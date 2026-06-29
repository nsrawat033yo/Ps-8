"""
LunaQuest — Hackathon-Winning Presentation Generator
=====================================================
Generates a 15-slide PowerPoint with dark theme, annotated maps,
speaker notes, and professional layout using python-pptx.
"""

import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ─── paths ───────────────────────────────────────────────────────────
BASE = r"c:\Users\nsraw\Downloads\tantra-mantra"
FIGS = os.path.join(BASE, "outputs", "figures")
OUT_PPT = os.path.join(BASE, "LunaQuest_Hackathon_Presentation.pptx")

# ─── colour palette ────────────────────────────────────────────────
BG_DARK   = RGBColor(0x0B, 0x0D, 0x17)   # deep space navy
BG_CARD   = RGBColor(0x12, 0x16, 0x26)   # card / content area
ACCENT_1  = RGBColor(0x00, 0xBF, 0xFF)   # cyan accent
ACCENT_2  = RGBColor(0xFF, 0x6B, 0x35)   # warm orange
ACCENT_3  = RGBColor(0x6C, 0x63, 0xFF)   # indigo/purple
WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT     = RGBColor(0xB0, 0xBE, 0xD0)   # light grey for body
GOLD      = RGBColor(0xFF, 0xD7, 0x00)
RED       = RGBColor(0xFF, 0x45, 0x45)
GREEN     = RGBColor(0x00, 0xE6, 0x76)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

prs = Presentation()
prs.slide_width  = SLIDE_W
prs.slide_height = SLIDE_H

# ═══════════════════════════════════════════════════════════════════
#  HELPER FUNCTIONS
# ═══════════════════════════════════════════════════════════════════

def set_slide_bg(slide, color=BG_DARK):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_text_box(slide, left, top, width, height, text,
                 font_size=18, color=WHITE, bold=False,
                 alignment=PP_ALIGN.LEFT, font_name="Calibri",
                 anchor=MSO_ANCHOR.TOP):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.auto_size = None
    try:
        tf.vertical_anchor = anchor
    except Exception:
        pass
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox


def add_multi_text(slide, left, top, width, height, lines,
                   font_size=16, color=LIGHT, font_name="Calibri",
                   line_space=1.3, bullet_color=ACCENT_1):
    """Add multiple lines / bullets."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.auto_size = None
    for i, line in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        # bullet icon
        if line.startswith("•") or line.startswith("-") or line.startswith("✦"):
            p.text = line
        else:
            p.text = line
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = font_name
        p.space_after = Pt(4)
        try:
            p.line_spacing = Pt(int(font_size * line_space))
        except Exception:
            pass
    return txBox


def add_accent_bar(slide, left, top, width=Inches(0.08), height=Inches(0.8),
                   color=ACCENT_1):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def add_image_safe(slide, img_name, left, top, width=None, height=None):
    path = os.path.join(FIGS, img_name)
    if not os.path.exists(path):
        # placeholder box
        add_text_box(slide, left, top, Inches(4), Inches(0.5),
                     f"[Image: {img_name}]", font_size=12, color=RED)
        return None
    kwargs = {}
    if width:  kwargs['width'] = width
    if height: kwargs['height'] = height
    return slide.shapes.add_picture(path, left, top, **kwargs)


def add_notes(slide, text):
    notes_slide = slide.notes_slide
    notes_slide.notes_text_frame.text = text


def add_card_bg(slide, left, top, width, height, color=BG_CARD, alpha=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def add_section_title(slide, title, subtitle=None):
    add_accent_bar(slide, Inches(0.6), Inches(0.5), height=Inches(0.7))
    add_text_box(slide, Inches(0.85), Inches(0.45), Inches(10), Inches(0.8),
                 title, font_size=32, color=WHITE, bold=True)
    if subtitle:
        add_text_box(slide, Inches(0.85), Inches(1.15), Inches(10), Inches(0.5),
                     subtitle, font_size=16, color=LIGHT)


def add_bottom_tag(slide, text="LunaQuest  |  Lunar Ice Candidate Screening & Mission Planning  |  Team Tantra-Mantra"):
    add_text_box(slide, Inches(0.5), Inches(7.05), Inches(12), Inches(0.35),
                 text, font_size=10, color=RGBColor(0x55, 0x65, 0x75),
                 alignment=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════════════
#  SLIDE 1 — TITLE
# ═══════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
set_slide_bg(slide)

# gradient-like strip at top
strip = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0), SLIDE_W, Inches(0.06))
strip.fill.solid()
strip.fill.fore_color.rgb = ACCENT_1
strip.line.fill.background()

add_text_box(slide, Inches(1.5), Inches(1.6), Inches(10), Inches(1.2),
             "LUNAQUEST", font_size=54, color=ACCENT_1, bold=True,
             alignment=PP_ALIGN.CENTER)

add_text_box(slide, Inches(1.5), Inches(2.7), Inches(10), Inches(0.8),
             "Lunar Subsurface Ice Candidate Screening & Autonomous Mission Planning",
             font_size=22, color=WHITE, bold=False,
             alignment=PP_ALIGN.CENTER)

add_text_box(slide, Inches(1.5), Inches(3.6), Inches(10), Inches(0.5),
             "Faustini Crater AOI  •  Lunar South Polar Region",
             font_size=16, color=LIGHT, alignment=PP_ALIGN.CENTER)

# team info card
add_card_bg(slide, Inches(4), Inches(4.6), Inches(5.3), Inches(1.5))
add_text_box(slide, Inches(4.2), Inches(4.7), Inches(5), Inches(0.4),
             "Team Tantra-Mantra", font_size=20, color=GOLD, bold=True,
             alignment=PP_ALIGN.CENTER)
add_text_box(slide, Inches(4.2), Inches(5.2), Inches(5), Inches(0.7),
             "Bharat Antariksh Hackathon 2026\nPowered by Chandrayaan-2 SAR / DFSAR  •  TMC-2 DTM  •  OHRC",
             font_size=14, color=LIGHT, alignment=PP_ALIGN.CENTER)

# bottom strip
strip2 = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
         Inches(0), Inches(7.44), SLIDE_W, Inches(0.06))
strip2.fill.solid()
strip2.fill.fore_color.rgb = ACCENT_1
strip2.line.fill.background()

add_notes(slide, """SPEAKER NOTES — Title Slide
Good [morning/afternoon], judges.

We are Team Tantra-Mantra, and today we present LunaQuest — a complete end-to-end pipeline for Lunar Subsurface Ice Candidate Screening and Autonomous Mission Planning.

This system uses real Chandrayaan-2 SAR/DFSAR radar data and TMC-2 Digital Terrain Models, focused on the Faustini Crater at the Lunar South Pole.

We want to emphasize from the outset — this system does NOT claim to confirm ice detection. What it does is provide a rigorous, explainable, and proxy-based candidate screening and mission planning framework.

Let us walk you through the full pipeline.""")

# ═══════════════════════════════════════════════════════════════════
#  SLIDE 2 — PROBLEM STATEMENT
# ═══════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_section_title(slide, "Why Lunar Ice Matters",
                  "The Foundation for Sustained Human Presence Beyond Earth")

bullets = [
    "✦  Water ice is the most critical in-situ resource for sustained lunar operations",
    "✦  Electrolysis splits H₂O → H₂ (fuel) + O₂ (life support & oxidizer)",
    "✦  1 kg of lunar water saves ~$1M in Earth-launch costs (NASA estimates)",
    "✦  Enables deep-space missions: Moon as a refueling depot for Mars & beyond",
    "✦  Chandrayaan-1 M³ & Mini-SAR first detected polar ice signatures (2009)",
    "✦  Chandrayaan-2 DFSAR provides the most detailed S-band radar data of the south pole",
]
add_multi_text(slide, Inches(0.85), Inches(1.8), Inches(7), Inches(4.5),
               bullets, font_size=18, color=LIGHT)

# Right-side highlight card
add_card_bg(slide, Inches(8.5), Inches(2.0), Inches(4.2), Inches(3.2), color=RGBColor(0x0A, 0x1A, 0x35))
add_text_box(slide, Inches(8.7), Inches(2.1), Inches(3.8), Inches(0.5),
             "ISRU Potential", font_size=20, color=ACCENT_1, bold=True,
             alignment=PP_ALIGN.CENTER)
isru_lines = [
    "🚀  Rocket propellant production",
    "💧  Drinking water for crew",
    "🌬️  Breathable oxygen generation",
    "🔬  Radiation shielding material",
    "🏗️  Construction: concrete from regolith + water",
]
add_multi_text(slide, Inches(8.7), Inches(2.7), Inches(3.8), Inches(2.5),
               isru_lines, font_size=15, color=LIGHT)

add_bottom_tag(slide)
add_notes(slide, """SPEAKER NOTES — Problem Statement
Why does lunar ice matter? It is the single most important in-situ resource for establishing a sustained human presence on the Moon.

Water ice, if present, can be split via electrolysis into hydrogen for rocket fuel and oxygen for life support. NASA estimates that every kilogram of water we don't have to launch from Earth saves approximately one million dollars.

Moreover, the Moon could serve as a refueling station — enabling deep-space missions to Mars and beyond.

Chandrayaan-1's Mini-SAR and M³ instruments first detected signatures consistent with polar ice in 2009. Chandrayaan-2's DFSAR now provides the highest-resolution S-band radar data ever collected at the lunar south pole, giving us an unprecedented opportunity to screen for subsurface ice candidates.

This is why our pipeline is built on Chandrayaan-2 data — it's India's best instrument for this task.""")

# ═══════════════════════════════════════════════════════════════════
#  SLIDE 3 — CHALLENGES
# ═══════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_section_title(slide, "Challenges We Address",
                  "Why Detecting Lunar Subsurface Ice Is One of the Hardest Problems in Planetary Science")

# Three challenge cards
challenges = [
    ("Data Gaps & Ambiguity", ACCENT_1, [
        "• No calibrated CPR/DOP products available",
        "• OHRC does not overlap Faustini AOI",
        "• No Diviner thermal / PSR illumination",
        "• Rough surfaces mimic ice radar signatures",
    ]),
    ("Extreme Terrain", ACCENT_2, [
        "• Permanently Shadowed Regions (PSRs)",
        "• Temperatures below 40 K in craters",
        "• Slopes > 15° create landing hazards",
        "• No GPS, no direct communication",
    ]),
    ("Uncertainty Everywhere", ACCENT_3, [
        "• High CPR ≠ ice (roughness ambiguity)",
        "• Proxy-based measurements only",
        "• Unknown ice depth & concentration",
        "• No ground truth available",
    ]),
]

for i, (title, accent, lines) in enumerate(challenges):
    x = Inches(0.5 + i * 4.2)
    add_card_bg(slide, x, Inches(2.0), Inches(3.8), Inches(4.5))
    add_accent_bar(slide, x + Inches(0.15), Inches(2.2), width=Inches(3.5),
                   height=Inches(0.06), color=accent)
    add_text_box(slide, x + Inches(0.2), Inches(2.4), Inches(3.4), Inches(0.5),
                 title, font_size=20, color=accent, bold=True)
    add_multi_text(slide, x + Inches(0.2), Inches(3.1), Inches(3.4), Inches(3.0),
                   lines, font_size=15, color=LIGHT)

add_bottom_tag(slide)
add_notes(slide, """SPEAKER NOTES — Challenges
Let me outline three critical challenges that make this problem extremely hard.

First — Data Gaps. There are no calibrated Circular Polarization Ratio products available for our AOI. The OHRC high-resolution camera data doesn't overlap our Faustini study area. We're missing Diviner thermal maps and real PSR illumination data. And rough surfaces produce high radar returns that look identical to ice signatures.

Second — Extreme Terrain. We're working at the lunar south pole, inside permanently shadowed craters where temperatures drop below 40 Kelvin. Slopes above 15 degrees create landing hazards. There's no GPS and limited communication windows.

Third — Uncertainty. A high CPR value does NOT confirm ice — it could be surface roughness, rock scattering, or volume scattering from ice. All our measurements are proxy-based, the actual ice depth and concentration are unknown, and there is zero ground truth to validate against.

Our system is designed to work within these constraints, not pretend they don't exist.""")

# ═══════════════════════════════════════════════════════════════════
#  SLIDE 4 — OUR APPROACH (Pipeline Overview)
# ═══════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_section_title(slide, "Our Approach: End-to-End Pipeline",
                  "From Raw Chandrayaan-2 Data → Candidate Screening → Landing → Rover → Resource Estimation")

# Pipeline steps as connected boxes
steps = [
    ("1", "Data Ingestion\n& Validation", ACCENT_1),
    ("2", "Ice Probability\nMapping", ACCENT_1),
    ("3", "Candidate Patch\nExtraction", ACCENT_2),
    ("4", "Landing Site\nSelection", ACCENT_2),
    ("5", "Rover Path\nPlanning", ACCENT_3),
    ("6", "Resource\nEstimation", ACCENT_3),
    ("7", "Explainable\nAI Layer", GOLD),
]

box_w = Inches(1.55)
box_h = Inches(1.2)
start_x = Inches(0.5)
y_top = Inches(2.2)

for i, (num, label, accent) in enumerate(steps):
    x = start_x + Inches(i * 1.78)
    # box
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
            x, y_top, box_w, box_h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = BG_CARD
    shape.line.color.rgb = accent
    shape.line.width = Pt(2)
    # number
    add_text_box(slide, x + Inches(0.05), y_top + Inches(0.05),
                 Inches(0.4), Inches(0.35),
                 num, font_size=16, color=accent, bold=True)
    # label
    add_text_box(slide, x + Inches(0.05), y_top + Inches(0.35),
                 box_w - Inches(0.1), Inches(0.75),
                 label, font_size=13, color=WHITE, alignment=PP_ALIGN.CENTER)
    # arrow
    if i < len(steps) - 1:
        add_text_box(slide, x + box_w - Inches(0.05), y_top + Inches(0.35),
                     Inches(0.35), Inches(0.35),
                     "→", font_size=22, color=ACCENT_1, bold=True)

# Key differentiation section
add_card_bg(slide, Inches(0.5), Inches(4.0), Inches(12.3), Inches(3.0))
add_text_box(slide, Inches(0.7), Inches(4.1), Inches(6), Inches(0.4),
             "KEY DIFFERENTIATORS", font_size=18, color=GOLD, bold=True)

diff_cols = [
    [
        "✦  AOI-aware data coverage validation",
        "✦  CPR-style ratio proxy with roughness penalty",
        "✦  Ambiguity-aware candidate filtering (57 → 53 refined)",
        "✦  Fuzzy multi-criteria landing site scoring",
    ],
    [
        "✦  4 A*-optimized rover route variants",
        "✦  Uncertainty-aware ice volume estimation",
        "✦  Explainable AI: 'why this patch? why this site?'",
        "✦  All limitations explicitly documented",
    ],
]
add_multi_text(slide, Inches(0.7), Inches(4.6), Inches(5.8), Inches(2.5),
               diff_cols[0], font_size=15, color=LIGHT)
add_multi_text(slide, Inches(6.8), Inches(4.6), Inches(5.8), Inches(2.5),
               diff_cols[1], font_size=15, color=LIGHT)

add_bottom_tag(slide)
add_notes(slide, """SPEAKER NOTES — Pipeline Overview
Our system is NOT a single model — it's a complete end-to-end pipeline with seven interconnected stages.

Stage 1: We ingest and validate 1,136 data products from Chandrayaan-2, selecting only those with verified AOI coverage. This is critical — many SAR and DTM products don't actually overlap our study area.

Stage 2: We compute an ice probability map using a weighted fusion of CPR-style ratio proxy (30%), SAR backscatter intensity (27%), roughness penalty (23%), and PSR stability proxy (20%).

Stage 3: From 57 initial detections, we refine to 53 candidates using ambiguity filtering that penalizes patches where roughness — not ice — likely explains the radar signature.

Stage 4: We select 5 landing candidates using fuzzy multi-criteria scoring: slope < 5°, low roughness, outside the ice candidate mask but close to science targets.

Stage 5: A* path planning produces 4 route variants — shortest, safest, energy-efficient, and science-priority — with the best route at approximately 8.95 km.

Stage 6: Probabilistic ice volume estimation using depth and concentration scenarios.

Stage 7: An explainability layer that breaks down exactly why each decision was made.

What makes us different: we explicitly handle data gaps, document all limitations, and never over-claim.""")

# ═══════════════════════════════════════════════════════════════════
#  SLIDE 5 — DATA SOURCES & COVERAGE
# ═══════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_section_title(slide, "Data Sources & Coverage Validation",
                  "AOI: Faustini F2  |  Lat: –87.8° to –86.9°  |  Lon: 80.0° to 85.0° E")

# Data source table
headers = ["Instrument", "Product Type", "Coverage", "Role", "Status"]
rows_data = [
    ["DFSAR/SAR", "SRI (LH/LV)", "100% AOI", "Ice candidate screening", "✅ Primary"],
    ["TMC-2", "DTM (d18)", "100% AOI", "Slope / roughness / landing", "✅ Active"],
    ["OHRC", "Calibrated IMG", "0% AOI", "Hazard / boulder review", "⚠️ No overlap"],
    ["Diviner", "Thermal map", "Not available", "Cold-trap validation", "❌ Future"],
    ["PSR/Illum.", "Shadow map", "Proxy only", "PSR stability context", "⚠️ Proxy"],
    ["LOLA", "Albedo / DEM", "Fallback", "Cross-validation", "✅ Context"],
]

# Header
y_start = Inches(2.0)
col_widths = [Inches(1.6), Inches(1.8), Inches(1.5), Inches(3.2), Inches(1.8)]
x_start = Inches(0.85)

for ci, h in enumerate(headers):
    x = x_start + sum(w for w in [Inches(0)] + col_widths[:ci])
    add_text_box(slide, x, y_start, col_widths[ci], Inches(0.4),
                 h, font_size=14, color=ACCENT_1, bold=True)

for ri, row in enumerate(rows_data):
    y = y_start + Inches(0.45 + ri * 0.5)
    row_bg = BG_CARD if ri % 2 == 0 else BG_DARK
    for ci, cell in enumerate(row):
        x = x_start + sum(w for w in [Inches(0)] + col_widths[:ci])
        clr = GREEN if "✅" in cell else (GOLD if "⚠️" in cell else (RED if "❌" in cell else LIGHT))
        add_text_box(slide, x, y, col_widths[ci], Inches(0.4),
                     cell, font_size=13, color=clr)

# OHRC handling callout
add_card_bg(slide, Inches(0.85), Inches(5.2), Inches(11.5), Inches(1.5),
            color=RGBColor(0x1A, 0x10, 0x05))
add_text_box(slide, Inches(1.0), Inches(5.3), Inches(11), Inches(0.4),
             "⚠️  OHRC Gap Handling — Explicit, Not Ignored", font_size=16, color=ACCENT_2, bold=True)
add_multi_text(slide, Inches(1.0), Inches(5.7), Inches(11), Inches(0.9), [
    "Available OHRC footprints cluster near –89° to –90° latitude — our AOI is at –86.9° to –87.8°",
    "System explicitly detects zero overlap and excludes OHRC from hazard scoring, rather than using invalid data",
    "This is documented as a future validation requirement — not hidden",
], font_size=13, color=LIGHT)

add_bottom_tag(slide)
add_notes(slide, """SPEAKER NOTES — Data Sources
We processed 1,136 data products from the Chandrayaan-2 data archive. Let me walk you through what we have and what we don't.

We have full 100% AOI coverage from DFSAR SAR — specifically the SRI intensity products with both LH and LV polarization channels. This is our primary screening instrument.

We also have 100% coverage from TMC-2 Digital Terrain Model, giving us slope and roughness information critical for landing and route planning.

Now — and this is important — we do NOT have OHRC data overlapping our AOI. The available OHRC footprints are near the pole at 89-90 degrees south, while our Faustini study area is at 86.9 to 87.8 degrees south. Rather than pretending this data exists or using it incorrectly, our system explicitly detects zero overlap and excludes OHRC from hazard scoring.

We're also missing Diviner thermal maps and real PSR illumination data. Our system uses a latitude-based PSR stability proxy as a placeholder, clearly marked as such.

This transparent handling of data gaps is a key strength of our system — we show exactly what data we have, what we're missing, and how it affects our confidence.""")

# ═══════════════════════════════════════════════════════════════════
#  SLIDE 6 — ICE PROBABILITY MAPPING
# ═══════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_section_title(slide, "Ice Probability Mapping",
                  "Weighted Fusion Model — CPR Proxy + Backscatter + Roughness + PSR Stability")

# Left: formula & weights
add_card_bg(slide, Inches(0.5), Inches(2.0), Inches(5.0), Inches(4.8))
add_text_box(slide, Inches(0.7), Inches(2.1), Inches(4.5), Inches(0.4),
             "Weighted Fusion Formula", font_size=18, color=ACCENT_1, bold=True)
add_text_box(slide, Inches(0.7), Inches(2.6), Inches(4.5), Inches(0.5),
             "P(ice) = w₁·CPR_proxy + w₂·SAR_intensity + w₃·(1-roughness) + w₄·PSR_proxy",
             font_size=13, color=GOLD)

weights = [
    ("CPR-Style Ratio Proxy (LH/LV)", "0.30", "Elevated ratio suggests volume scattering"),
    ("SAR Backscatter Intensity", "0.27", "Strong returns from subsurface interfaces"),
    ("Roughness Penalty", "0.23", "High roughness → reduces confidence"),
    ("PSR Stability Proxy", "0.20", "Latitude-based persistent shadow estimate"),
]
y = Inches(3.3)
for label, weight, desc in weights:
    add_text_box(slide, Inches(0.7), y, Inches(3.8), Inches(0.3),
                 f"▸  {label}: {weight}", font_size=14, color=WHITE, bold=True)
    add_text_box(slide, Inches(1.0), y + Inches(0.3), Inches(3.5), Inches(0.3),
                 desc, font_size=12, color=LIGHT)
    y += Inches(0.6)

# right: map
add_image_safe(slide, "ice_probability_map_annotated.png",
               Inches(5.8), Inches(1.8), width=Inches(7.0), height=Inches(5.2))

add_bottom_tag(slide)
add_notes(slide, """SPEAKER NOTES — Ice Probability Mapping
This is our core ice probability model. We use a weighted fusion approach with four components.

CPR-style Ratio Proxy at 30% weight: We compute the LH-to-LV ratio from the SRI products. I want to be clear — this is NOT calibrated CPR. The available products are real-valued intensity rasters, not the complex or Stokes products needed for true CPR/DOP computation. We call it a 'CPR-style ratio proxy' and it is used as a screening indicator only.

SAR Backscatter Intensity at 27%: Strong radar returns can indicate subsurface volume scattering from ice-regolith mixtures.

Roughness Penalty at 23%: This is crucial. High radar returns can also come from rough surfaces. Our roughness penalty reduces the ice probability for areas with high surface roughness, directly addressing the CPR-roughness ambiguity documented in the literature.

PSR Stability Proxy at 20%: Since we don't have real illumination data, we use a latitude-based proxy for persistent shadow stability. This is a placeholder — not a substitute for real PSR modeling.

On the right, you can see the resulting ice probability map over the Faustini AOI. The warm colors indicate higher probability regions.""")

# ═══════════════════════════════════════════════════════════════════
#  SLIDE 7 — CANDIDATE PATCH DETECTION
# ═══════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_section_title(slide, "Candidate Patch Detection & Filtering",
                  "57 Initial Detections → 53 Refined Candidates  |  Ambiguity-Aware Filtering")

# Left: stats
add_card_bg(slide, Inches(0.5), Inches(2.0), Inches(5.3), Inches(2.3))
add_text_box(slide, Inches(0.7), Inches(2.1), Inches(4.8), Inches(0.4),
             "Detection Statistics", font_size=18, color=ACCENT_1, bold=True)
stats_lines = [
    "▸  Initial candidate patches: 57",
    "▸  After ambiguity filtering: 53 refined",
    "▸  Candidate mask area: 0.375% of valid AOI",
    "▸  Top candidate: C-025 (13,125 m²  •  Score: 0.846  •  High confidence)",
    "▸  Confidence levels: High / Medium / Low",
]
add_multi_text(slide, Inches(0.7), Inches(2.6), Inches(4.8), Inches(1.7),
               stats_lines, font_size=14, color=LIGHT)

# top candidates table
add_card_bg(slide, Inches(0.5), Inches(4.5), Inches(5.3), Inches(2.5))
add_text_box(slide, Inches(0.7), Inches(4.6), Inches(4.8), Inches(0.4),
             "Top Ranked Candidate Patches", font_size=16, color=GOLD, bold=True)
top_cands = [
    "C-025  │  13,125 m²  │  Score 0.846  │  HIGH  │  Rank 1",
    "C-042  │   9,375 m²  │  Score 0.800  │  HIGH  │  Rank 2",
    "C-011  │  11,250 m²  │  Score 0.780  │  MED   │  Rank 3",
    "C-022  │   8,750 m²  │  Score 0.834  │  MED   │  Rank 4",
    "C-055  │   5,000 m²  │  Score 0.823  │  MED   │  Rank 5",
]
add_multi_text(slide, Inches(0.7), Inches(5.1), Inches(4.8), Inches(1.8),
               top_cands, font_size=12, color=LIGHT)

# Right: map
add_image_safe(slide, "ice_candidate_detection_map.png",
               Inches(6.2), Inches(1.8), width=Inches(6.8), height=Inches(5.3))

add_bottom_tag(slide)
add_notes(slide, """SPEAKER NOTES — Candidate Patch Detection
From the ice probability map, we extract connected candidate patches using threshold-based screening on multiple SAR features simultaneously — ratio, intensity, texture, and composite score.

We initially detect 57 candidate patches. After applying ambiguity filtering — which penalizes patches where high roughness likely explains the radar signature rather than ice — we refine this to 53 candidates. The filtered patches are not deleted; they're flagged for future validation.

Our top candidate, C-025, covers 13,125 square meters with a mean score of 0.846 and is classified as High confidence. The total candidate area is only 0.375% of the valid AOI — we are being conservative and selective.

Each patch is characterized by area, equivalent diameter, mean score, roughness ambiguity risk, slope context, threshold stability across different thresholds, and proximity to landing candidates.

The confidence classification combines score strength, patch extent, stability, uncertainty, and terrain context. This is NOT a simple threshold — it's a multi-factor assessment.""")

# ═══════════════════════════════════════════════════════════════════
#  SLIDE 8 — TECHNOLOGIES USED
# ═══════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_section_title(slide, "Technologies Used in Solution",
                  "Open-Source Stack Built for Scientific Reproducibility & Scalability")

# Technology categories in cards
tech_cats = [
    ("Remote Sensing & GIS", ACCENT_1, [
        "▸ GDAL / Rasterio — Geospatial raster I/O",
        "▸ Shapely / GeoPandas — Vector analysis",
        "▸ Pyproj — Coordinate transforms",
        "▸ QGIS (visual validation)",
    ]),
    ("Scientific Computing", ACCENT_2, [
        "▸ NumPy / SciPy — Array computation",
        "▸ scikit-image — Connected component analysis",
        "▸ Pandas — Tabular data processing",
        "▸ Matplotlib — Visualization & maps",
    ]),
    ("Machine Learning", ACCENT_3, [
        "▸ PyTorch — U-Net segmentation model",
        "▸ Custom TinyU-Net (5-channel input)",
        "▸ Weakly supervised pseudo-label training",
        "▸ Data augmentation (flip/rotate/noise)",
    ]),
    ("Path Planning & Optimization", GOLD, [
        "▸ A* algorithm — Multi-objective routing",
        "▸ Custom cost functions (slope/energy/risk)",
        "▸ Fuzzy logic — Landing site scoring",
        "▸ Connected-component patch extraction",
    ]),
]

for i, (title, accent, items) in enumerate(tech_cats):
    row = i // 2
    col = i % 2
    x = Inches(0.5 + col * 6.3)
    y = Inches(1.9 + row * 2.7)
    add_card_bg(slide, x, y, Inches(5.9), Inches(2.4))
    add_accent_bar(slide, x + Inches(0.1), y + Inches(0.15),
                   width=Inches(5.7), height=Inches(0.05), color=accent)
    add_text_box(slide, x + Inches(0.2), y + Inches(0.3), Inches(5.5), Inches(0.4),
                 title, font_size=18, color=accent, bold=True)
    add_multi_text(slide, x + Inches(0.2), y + Inches(0.8), Inches(5.5), Inches(1.5),
                   items, font_size=14, color=LIGHT)

# Language badge
add_card_bg(slide, Inches(10.5), Inches(6.6), Inches(2.5), Inches(0.5), color=RGBColor(0x0A, 0x1A, 0x35))
add_text_box(slide, Inches(10.5), Inches(6.6), Inches(2.5), Inches(0.5),
             "🐍  Python 3.x  |  100% Open Source", font_size=11, color=GREEN,
             alignment=PP_ALIGN.CENTER)

add_bottom_tag(slide)
add_notes(slide, """SPEAKER NOTES — Technologies Used
Our entire solution is built on a 100% open-source Python stack, ensuring reproducibility and transparency.

For Remote Sensing, we use GDAL and Rasterio for geospatial raster I/O, Shapely and GeoPandas for vector analysis, and Pyproj for coordinate transformations between lunar coordinate systems.

For Scientific Computing, NumPy and SciPy handle our core array computations, scikit-image performs connected component analysis for candidate patch extraction, and Matplotlib generates all our publication-quality maps and charts.

For Machine Learning, we built a custom TinyU-Net in PyTorch with 5-channel input — SAR intensity, ratio proxy, texture, polarization imbalance, and candidate score. This is trained with weak supervision against rule-based pseudo-labels.

For Path Planning, we implemented A* algorithm with custom multi-objective cost functions that balance slope penalty, energy consumption, traverse risk, and science reward. Landing site scoring uses fuzzy membership functions inspired by published research.

Everything runs in a single configurable JSON pipeline — no manual steps, fully reproducible.""")

# ═══════════════════════════════════════════════════════════════════
#  SLIDE 9 — ESTIMATED IMPLEMENTATION COST
# ═══════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_section_title(slide, "Estimated Implementation Cost",
                  "From Prototype to Operational Mission Planning Tool")

# Cost categories
costs = [
    ("Development & Integration", "₹15–20 Lakhs", ACCENT_1, [
        "▸ Software engineering (3–4 months, 3 developers)",
        "▸ Pipeline hardening, unit testing & CI/CD",
        "▸ Integration with ISRO PRADAN / ISSDC data APIs",
        "▸ Documentation & scientific validation reports",
    ]),
    ("Compute Infrastructure", "₹5–8 Lakhs", ACCENT_2, [
        "▸ GPU cloud instances for U-Net training at scale",
        "▸ High-memory nodes for full-resolution raster processing",
        "▸ Storage for multi-terabyte Chandrayaan-2 archive",
        "▸ Annual cloud / HPC maintenance",
    ]),
    ("Data Acquisition & Validation", "₹3–5 Lakhs", ACCENT_3, [
        "▸ Diviner thermal, LOLA albedo, PSR data licensing/download",
        "▸ Calibrated OHRC acquisition for Faustini overlap",
        "▸ Complex/Stokes SAR product procurement",
        "▸ Manual expert review of candidate patches",
    ]),
    ("Scaling & Operations", "₹5–8 Lakhs / year", GOLD, [
        "▸ Extend to 10+ AOIs across lunar south pole",
        "▸ Real-time data ingestion from ongoing missions",
        "▸ Dashboard / visualization for mission planners",
        "▸ Collaboration interface for ISRO scientists",
    ]),
]

for i, (title, cost, accent, items) in enumerate(costs):
    row = i // 2
    col = i % 2
    x = Inches(0.5 + col * 6.3)
    y = Inches(1.9 + row * 2.7)
    add_card_bg(slide, x, y, Inches(5.9), Inches(2.4))
    add_text_box(slide, x + Inches(0.2), y + Inches(0.15), Inches(3.5), Inches(0.4),
                 title, font_size=16, color=accent, bold=True)
    add_text_box(slide, x + Inches(3.8), y + Inches(0.15), Inches(1.9), Inches(0.4),
                 cost, font_size=16, color=GOLD, bold=True, alignment=PP_ALIGN.RIGHT)
    add_multi_text(slide, x + Inches(0.2), y + Inches(0.7), Inches(5.5), Inches(1.5),
                   items, font_size=13, color=LIGHT)

# Total bar
add_card_bg(slide, Inches(3.5), Inches(7.0), Inches(6.3), Inches(0.4),
            color=RGBColor(0x0A, 0x2A, 0x15))
add_text_box(slide, Inches(3.5), Inches(7.0), Inches(6.3), Inches(0.4),
             "Total Estimated Cost:  ₹28–41 Lakhs (Year 1)  |  ₹5–8 Lakhs/year thereafter",
             font_size=14, color=GREEN, bold=True, alignment=PP_ALIGN.CENTER)

add_bottom_tag(slide)
add_notes(slide, """SPEAKER NOTES — Implementation Cost
Let me walk you through the realistic cost estimates for taking this prototype to an operational tool.

Development and Integration: 15 to 20 lakhs covers 3-4 months of software engineering with 3 developers. This includes hardening the pipeline with proper testing, integrating with ISRO's PRADAN data portal, and generating scientific validation documentation.

Compute Infrastructure: 5 to 8 lakhs for GPU cloud instances for U-Net training at scale, high-memory processing for full-resolution rasters — our current SAR products are at 25-meter resolution and the high-res ones are at 4-meter resolution.

Data Acquisition: 3 to 5 lakhs for obtaining the missing validation datasets — Diviner thermal, LOLA albedo, real PSR maps, and calibrated OHRC overlapping our AOI. This also covers expert scientist time for manual review.

Scaling and Operations: 5 to 8 lakhs per year to extend coverage to 10+ AOIs across the entire lunar south pole, build a real-time ingestion pipeline, and create a dashboard for ISRO mission planners.

The total Year 1 cost is 28 to 41 lakhs — a fraction of any actual mission cost, but it delivers a validated mission planning tool that could directly support site selection for India's future lunar missions.""")

# ═══════════════════════════════════════════════════════════════════
#  SLIDE 10 — LANDING SITE SELECTION
# ═══════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_section_title(slide, "Landing Site Selection",
                  "Fuzzy Multi-Criteria Scoring  |  5 Candidates Identified  |  L-01 Recommended")

# L-01 details card
add_card_bg(slide, Inches(0.5), Inches(2.0), Inches(5.5), Inches(5.0))
add_text_box(slide, Inches(0.7), Inches(2.1), Inches(4.5), Inches(0.4),
             "🏆 Recommended: L-01", font_size=22, color=GOLD, bold=True)

l01_details = [
    "▸  Lat: –87.147°  |  Lon: 80.103° E",
    "▸  Suitability score: 0.920 / 1.0",
    "▸  Local slope: 0.0° (flat terrain — perfect)",
    "▸  Roughness hazard score: 1.0 (very safe)",
    "▸  Distance to nearest ice candidate: 75 m",
    "",
    "WHY L-01?",
    "▸  ✅  Outside radar candidate mask (no risk of",
    "       contaminating science target)",
    "▸  ✅  Below 5° slope cutoff for safe landing",
    "▸  ✅  Within 75 m of C-001 for quick rover access",
    "▸  ✅  Inside configured AOI boundary",
    "▸  ⚠️  Illumination & thermal: neutral placeholders",
    "       (real layers needed for final validation)",
]
add_multi_text(slide, Inches(0.7), Inches(2.7), Inches(5.0), Inches(4.0),
               l01_details, font_size=14, color=LIGHT)

# right: map
add_image_safe(slide, "landing_site_map_annotated.png",
               Inches(6.3), Inches(1.8), width=Inches(6.6), height=Inches(5.3))

add_bottom_tag(slide)
add_notes(slide, """SPEAKER NOTES — Landing Site Selection
We identified 5 preliminary landing candidates using a fuzzy multi-criteria scoring model inspired by published research on lunar south pole landing site selection.

Our scoring considers: slope safety, candidate proximity, roughness avoidance, illumination, and thermal factors. Since illumination and thermal data are not available, those components use neutral 0.5 placeholders — clearly marked as such.

L-01 is our recommended site with a suitability score of 0.920 out of 1.0. It sits at 87.147 degrees south, 80.103 degrees east, on perfectly flat terrain with zero slope. It's located just 75 meters from the nearest ice candidate patch but — critically — it's outside the candidate mask. This matters because you don't want to land ON your science target and contaminate it.

I want to emphasize: these are preliminary landing candidates, NOT certified landing products. Real illumination, thermal, communication line-of-sight, and boulder-level hazard assessment are all required before any site can be certified for an actual mission.""")

# ═══════════════════════════════════════════════════════════════════
#  SLIDE 11 — ROVER PATH PLANNING
# ═══════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_section_title(slide, "Rover Path Planning",
                  "A* Multi-Objective Optimization  |  4 Route Variants  |  L-01 → C-055")

# Route comparison
add_card_bg(slide, Inches(0.5), Inches(2.0), Inches(5.8), Inches(2.8))
add_text_box(slide, Inches(0.7), Inches(2.1), Inches(5.4), Inches(0.4),
             "Route Comparison Table", font_size=16, color=ACCENT_1, bold=True)

routes_header = "Route Type       │ Length   │ Energy  │ Risk   │ Science │ Status"
routes_data = [
    "🥇 Shortest        │ 8,949 m  │ 10,319  │ 0.062  │ 0.583   │ ★ Recommended",
    "🛡️ Safest           │ 8,949 m  │ 10,319  │ 0.062  │ 0.583   │ Alternative",
    "🔬 Science Priority │ 8,949 m  │ 10,494  │ 0.070  │ 0.675   │ Alternative",
    "⚡ Energy Efficient  │ 8,949 m  │ 10,319  │ 0.062  │ 0.583   │ Alternative",
]
add_text_box(slide, Inches(0.7), Inches(2.5), Inches(5.4), Inches(0.3),
             routes_header, font_size=11, color=GOLD, bold=True)
add_multi_text(slide, Inches(0.7), Inches(2.9), Inches(5.4), Inches(1.6),
               routes_data, font_size=12, color=LIGHT)

# Simulation stats
add_card_bg(slide, Inches(0.5), Inches(5.0), Inches(5.8), Inches(2.0))
add_text_box(slide, Inches(0.7), Inches(5.1), Inches(5.4), Inches(0.4),
             "Traversal Simulation", font_size=16, color=ACCENT_2, bold=True)
sim_lines = [
    "▸  Simulated steps: 296 waypoints",
    "▸  100% of route under 5° slope (safe terrain)",
    "▸  No blocked cells (> 15° slope) encountered",
    "▸  Cumulative energy proxy: 9,051 units",
    "▸  Traverse risk score: 0.062 (very low)",
]
add_multi_text(slide, Inches(0.7), Inches(5.5), Inches(5.4), Inches(1.4),
               sim_lines, font_size=14, color=LIGHT)

# right: route map
add_image_safe(slide, "rover_routes_annotated.png",
               Inches(6.6), Inches(1.8), width=Inches(6.4), height=Inches(5.3))

add_bottom_tag(slide)
add_notes(slide, """SPEAKER NOTES — Rover Path Planning
We use A-star path planning with four different cost function variants to generate route options from landing site L-01 to science target C-055.

The cost functions balance:
- Slope penalty for terrain safety
- Energy consumption proxy based on elevation changes
- Traverse risk from steep sections
- Science reward based on proximity to additional candidate patches along the route

All four routes converge to approximately 8,949 meters — about 8.95 km — with 100% of each route traversing terrain under 5 degrees slope. No blocked cells above 15 degrees were encountered.

The Science Priority route has a slightly higher science reward score of 0.675 because it routes closer to additional ice candidate patches along the way, at the cost of slightly higher energy.

We recommend the Shortest route as it achieves the best balance of distance, energy, and risk, with the lowest traverse risk score of 0.062.

The traversal simulation ran 296 waypoints with cumulative energy proxy of about 9,051 units. Remember — these are conceptual planning routes, NOT operational rover commands.""")

# ═══════════════════════════════════════════════════════════════════
#  SLIDE 12 — ICE VOLUME ESTIMATION
# ═══════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_section_title(slide, "Ice Volume Estimation (Innovation)",
                  "Scenario-Based Probabilistic Resource Assessment  |  Planning-Only Estimates")

# Formula
add_card_bg(slide, Inches(0.5), Inches(2.0), Inches(6.2), Inches(1.3))
add_text_box(slide, Inches(0.7), Inches(2.1), Inches(5.8), Inches(0.4),
             "Volume Model", font_size=16, color=ACCENT_1, bold=True)
add_text_box(slide, Inches(0.7), Inches(2.5), Inches(5.8), Inches(0.3),
             "V_ice = Σ (patch_area × assumed_depth × ice_fraction)", font_size=14, color=GOLD)
add_text_box(slide, Inches(0.7), Inches(2.9), Inches(5.8), Inches(0.3),
             "Mass = V_ice × 917 kg/m³ (ice density)", font_size=12, color=LIGHT)

# Scenario table for C-025
add_card_bg(slide, Inches(0.5), Inches(3.5), Inches(6.2), Inches(3.5))
add_text_box(slide, Inches(0.7), Inches(3.6), Inches(5.8), Inches(0.4),
             "C-025 Resource Scenarios (13,125 m²  |  HIGH confidence)", font_size=15, color=GOLD, bold=True)

scenario_header = "Depth(m) │ Ice Frac. │ Volume(m³) │ Mass(kg)"
scenarios = [
    "1 m      │   1%      │    131      │  120,356",
    "1 m      │   3%      │    394      │  361,069",
    "1 m      │  10%      │  1,313      │ 1,203,563",
    "3 m      │   3%      │  1,181      │ 1,083,206",
    "3 m      │  10%      │  3,938      │ 3,610,688",
    "5 m      │  10%      │  6,563      │ 6,017,813",
]
add_text_box(slide, Inches(0.7), Inches(4.0), Inches(5.8), Inches(0.3),
             scenario_header, font_size=12, color=ACCENT_1, bold=True)
add_multi_text(slide, Inches(0.7), Inches(4.4), Inches(5.8), Inches(2.5),
               scenarios, font_size=12, color=LIGHT)

# Depth likelihood card
add_card_bg(slide, Inches(7.0), Inches(2.0), Inches(5.8), Inches(2.0))
add_text_box(slide, Inches(7.2), Inches(2.1), Inches(5.4), Inches(0.4),
             "Depth Likelihood Analysis", font_size=16, color=ACCENT_2, bold=True)
depth_lines = [
    "▸  C-055: Shallow likelihood = 1.000",
    "▸  Deep likelihood = 0.650",
    "▸  Classification: 'shallow-likely'",
    "▸  Based on radar score + roughness context",
    "▸  NOT measured depth — planning labels only",
]
add_multi_text(slide, Inches(7.2), Inches(2.6), Inches(5.4), Inches(1.4),
               depth_lines, font_size=14, color=LIGHT)

# Depth map
add_image_safe(slide, "depth_likelihood_map.png",
               Inches(7.0), Inches(4.2), width=Inches(5.8), height=Inches(3.0))

# Warning
add_card_bg(slide, Inches(0.5), Inches(7.0), Inches(12.3), Inches(0.4),
            color=RGBColor(0x2A, 0x15, 0x05))
add_text_box(slide, Inches(0.7), Inches(7.0), Inches(11.8), Inches(0.4),
             "⚠️  These are scenario-based planning estimates only — NOT measured resource quantities",
             font_size=12, color=ACCENT_2, bold=True, alignment=PP_ALIGN.CENTER)

add_bottom_tag(slide)
add_notes(slide, """SPEAKER NOTES — Ice Volume Estimation
This is one of our key innovations — a scenario-based probabilistic resource estimation layer.

The model is straightforward: Volume = patch area times assumed depth times assumed ice fraction. Mass uses the standard ice density of 917 kg per cubic meter.

For our top candidate C-025, with an area of 13,125 square meters, the scenarios range from a conservative estimate of 131 cubic meters and 120,000 kg of ice (at 1 meter depth, 1% ice fraction) to an optimistic estimate of 6,563 cubic meters and over 6 million kg of ice (at 5 meters depth, 10% ice fraction).

We also compute depth likelihood classifications based on radar score and roughness — C-055 shows a shallow likelihood of 1.0 and deep likelihood of 0.65, classified as 'shallow-likely'. These are rule-based planning labels, not measured depths.

I want to stress strongly: these are NOT measured resource quantities. They are scenario-based planning estimates that only become meaningful after independent validation of the candidate patches. This is for mission planners to assess whether a rover mission to a particular candidate is worth the investment.""")

# ═══════════════════════════════════════════════════════════════════
#  SLIDE 13 — EXPLAINABILITY LAYER
# ═══════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_section_title(slide, "Explainable AI Layer",
                  "Every Decision Has a Traceable Justification  |  No Black Boxes")

# Three explainability cards
explain_items = [
    ("Why this Patch?\n(C-055 Example)", ACCENT_1, [
        "▸  Mean candidate probability: 0.839",
        "▸  CPR-proxy contribution: 0.281 (30%)",
        "▸  SAR intensity: 0.267 (27%)",
        "▸  Roughness suitability: 0.223 (23%)",
        "▸  PSR-proxy: 0.069 (20% × proxy value)",
        "▸  Roughness ambiguity risk: LOW",
        "▸  Ice confidence score: 1.000",
    ]),
    ("Why this Landing Site?\n(L-01 Example)", ACCENT_2, [
        "▸  Slope score: 1.0 (0° slope — flat)",
        "▸  Proximity score: 0.930 (75 m away)",
        "▸  Roughness score: 1.0 (safe terrain)",
        "▸  Outside candidate mask: YES ✅",
        "▸  Inside AOI: YES ✅",
        "▸  Not in steep zone: YES ✅",
        "▸  Illumination: placeholder (future need)",
    ]),
    ("Why this Route?\n(Shortest Route)", ACCENT_3, [
        "▸  Length: 8,949 m (≈ 8.95 km)",
        "▸  Energy proxy: 10,319 (lowest)",
        "▸  Traverse risk: 0.062 (very low)",
        "▸  100% under 5° slope",
        "▸  0 blocked cells avoided",
        "▸  Science reward: 0.583",
        "▸  Decision rank: #1 of 4 variants",
    ]),
]

for i, (title, accent, items) in enumerate(explain_items):
    x = Inches(0.3 + i * 4.3)
    add_card_bg(slide, x, Inches(2.0), Inches(4.0), Inches(4.8))
    add_accent_bar(slide, x + Inches(0.1), Inches(2.15),
                   width=Inches(3.8), height=Inches(0.05), color=accent)
    add_text_box(slide, x + Inches(0.2), Inches(2.35), Inches(3.6), Inches(0.7),
                 title, font_size=16, color=accent, bold=True)
    add_multi_text(slide, x + Inches(0.2), Inches(3.2), Inches(3.6), Inches(3.5),
                   items, font_size=13, color=LIGHT)

add_bottom_tag(slide)
add_notes(slide, """SPEAKER NOTES — Explainability Layer
A critical feature of LunaQuest is that every decision is fully explainable. There are no black boxes.

For patch selection, we can break down exactly how much each component contributes. For C-055: the CPR proxy contributes 0.281, SAR intensity contributes 0.267, roughness suitability contributes 0.223, and the PSR proxy contributes 0.069. We can see that roughness ambiguity risk is low, meaning the high radar return is more likely from ice-like scattering than surface roughness.

For landing site L-01, we show the complete scoring breakdown: perfect slope score, high proximity score, full roughness safety, and explicit documentation that illumination and thermal layers are placeholders.

For the route, we show length, energy, risk, slope percentages, and how it ranks among all variants.

This level of transparency is essential for judges, for mission planners, and for scientific credibility. If a decision seems wrong, anyone can trace back through the contribution factors and understand why.""")

# ═══════════════════════════════════════════════════════════════════
#  SLIDE 14 — RESULTS SUMMARY
# ═══════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_section_title(slide, "Results Summary",
                  "From Detection → Landing → Route → Resource Estimation")

# Key metrics in large cards
metrics = [
    ("53", "Refined Ice\nCandidates", ACCENT_1),
    ("C-025", "Top Patch\n(Score: 0.846)", ACCENT_2),
    ("L-01", "Best Landing\n(Score: 0.920)", ACCENT_3),
    ("8.95 km", "Best Route\n(100% safe)", GOLD),
    ("6M+ kg", "Max Resource\nScenario", GREEN),
]

for i, (val, label, accent) in enumerate(metrics):
    x = Inches(0.3 + i * 2.55)
    add_card_bg(slide, x, Inches(2.0), Inches(2.3), Inches(2.0))
    add_text_box(slide, x + Inches(0.1), Inches(2.15), Inches(2.1), Inches(0.7),
                 val, font_size=30, color=accent, bold=True,
                 alignment=PP_ALIGN.CENTER)
    add_text_box(slide, x + Inches(0.1), Inches(2.9), Inches(2.1), Inches(0.7),
                 label, font_size=13, color=LIGHT,
                 alignment=PP_ALIGN.CENTER)

# Impact statement
add_card_bg(slide, Inches(0.5), Inches(4.4), Inches(12.3), Inches(1.2),
            color=RGBColor(0x0A, 0x1A, 0x35))
add_text_box(slide, Inches(0.7), Inches(4.5), Inches(11.8), Inches(0.4),
             "IMPACT: Complete Pipeline — Detection → Landing → Route → Resource Assessment",
             font_size=20, color=GOLD, bold=True, alignment=PP_ALIGN.CENTER)
add_text_box(slide, Inches(0.7), Inches(4.95), Inches(11.8), Inches(0.5),
             "This is NOT a single-model solution. It is a full mission planning framework that can be extended and validated incrementally.",
             font_size=14, color=LIGHT, alignment=PP_ALIGN.CENTER)

# Combined decision map
add_image_safe(slide, "combined_decision_map.png",
               Inches(3.0), Inches(5.7), width=Inches(7.0), height=Inches(1.6))

add_bottom_tag(slide)
add_notes(slide, """SPEAKER NOTES — Results Summary
Let me summarize our key results.

53 refined ice candidate patches from 57 initial detections — filtered through ambiguity-aware screening.

C-025 is our top-ranked candidate with a score of 0.846, covering 13,125 square meters, classified as HIGH confidence.

L-01 is our recommended landing site with a suitability score of 0.920, located on flat terrain just 75 meters from the nearest ice candidate.

Our best rover route is approximately 8.95 kilometers, with 100% of the traverse under 5 degrees slope — extremely safe terrain.

And our maximum resource scenario estimates over 6 million kilograms of potential ice for the top candidate at optimistic depth and concentration assumptions.

The key impact here is that this is a COMPLETE pipeline. We go from raw data to candidate screening to landing to routing to resource estimation, with full explainability at every step. This is not a single model — it's a mission planning framework.""")

# ═══════════════════════════════════════════════════════════════════
#  SLIDE 15 — LIMITATIONS (Critical for credibility)
# ═══════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_section_title(slide, "Limitations & Scientific Boundaries",
                  "What We Can Claim vs. What Requires Further Validation")

# Two columns
add_card_bg(slide, Inches(0.5), Inches(2.0), Inches(6.0), Inches(4.8),
            color=RGBColor(0x1A, 0x0A, 0x0A))
add_text_box(slide, Inches(0.7), Inches(2.1), Inches(5.5), Inches(0.4),
             "❌  What We Do NOT Claim", font_size=18, color=RED, bold=True)
no_claims = [
    "✖  No direct compositional proof of ice",
    "✖  CPR-style proxy ≠ calibrated CPR/DOP",
    "✖  Candidate scores are screening outputs,\n     not confirmed ice detections",
    "✖  Landing sites are preliminary, not certified",
    "✖  Routes are conceptual, not operational",
    "✖  U-Net metrics = pseudo-label agreement only",
    "✖  Resource estimates are scenario-based,\n     not measured quantities",
]
add_multi_text(slide, Inches(0.7), Inches(2.7), Inches(5.5), Inches(3.5),
               no_claims, font_size=14, color=RGBColor(0xFF, 0x88, 0x88))

add_card_bg(slide, Inches(6.8), Inches(2.0), Inches(6.0), Inches(4.8),
            color=RGBColor(0x0A, 0x15, 0x0A))
add_text_box(slide, Inches(7.0), Inches(2.1), Inches(5.5), Inches(0.4),
             "✅  What We DO Provide", font_size=18, color=GREEN, bold=True)
yes_claims = [
    "✔  Rigorous proxy-based candidate screening",
    "✔  Full data coverage validation (AOI-aware)",
    "✔  Roughness ambiguity penalty (not ignored)",
    "✔  Multi-criteria landing site ranking",
    "✔  Multi-objective rover route planning",
    "✔  Probabilistic resource scenario assessment",
    "✔  Full explainability for every decision",
    "✔  Every limitation explicitly documented",
]
add_multi_text(slide, Inches(7.0), Inches(2.7), Inches(5.5), Inches(3.5),
               yes_claims, font_size=14, color=RGBColor(0x88, 0xFF, 0x88))

# Footer emphasis
add_text_box(slide, Inches(0.5), Inches(6.9), Inches(12.3), Inches(0.4),
             "\"Transparency about limitations is not a weakness — it is the foundation of scientific credibility.\"",
             font_size=14, color=GOLD, alignment=PP_ALIGN.CENTER, bold=True)

add_bottom_tag(slide)
add_notes(slide, """SPEAKER NOTES — Limitations
This is arguably our most important slide, because scientific credibility depends on honest communication of limitations.

What we do NOT claim:
- We have no direct compositional proof of ice. Period.
- Our CPR-style proxy is NOT calibrated CPR or DOP — the available products are real-valued SRI intensity, not complex or Stokes products.
- Candidate scores are screening outputs that require independent validation.
- Landing sites are preliminary — they need illumination, thermal, boulder, and communication assessment before certification.
- Rover routes are conceptual planning products, not operational commands.
- Our U-Net metrics measure agreement with rule-based pseudo-labels, not ground truth.
- Resource estimates are scenarios with assumed depth and ice fraction.

What we DO provide:
- A rigorous, reproducible screening pipeline built on real Chandrayaan-2 data.
- Full data coverage validation — we explicitly handle missing OHRC, missing Diviner, missing PSR.
- Roughness ambiguity is directly penalized — we don't ignore the biggest confound.
- Every decision is explainable with contribution breakdowns.
- And every limitation is documented, not hidden.

We believe this transparency makes our system MORE useful, not less.""")

# ═══════════════════════════════════════════════════════════════════
#  SLIDE 16 — FUTURE WORK
# ═══════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_section_title(slide, "Future Work & Roadmap",
                  "From Prototype to Production-Ready Mission Planning Tool")

# Phase cards
phases = [
    ("Phase 1: Validation (3 months)", ACCENT_1, [
        "▸  Download calibrated OHRC overlapping Faustini AOI",
        "▸  Integrate Diviner thermal maps (< 110 K cold-trap screening)",
        "▸  Add real PSR / illumination persistence layer",
        "▸  Replace CPR proxy with complex/Stokes-derived CPR/DOP",
    ]),
    ("Phase 2: Enhancement (6 months)", ACCENT_2, [
        "▸  Multi-pass SAR temporal consistency analysis",
        "▸  LOLA albedo + LAMP UV H₂O proxy integration",
        "▸  M³ spectral ice signature cross-validation",
        "▸  Dynamic illumination / communication in rover planning",
    ]),
    ("Phase 3: Scaling (12 months)", ACCENT_3, [
        "▸  Extend to 10+ AOIs across lunar south pole",
        "▸  Real-time ISRO data ingestion pipeline",
        "▸  Interactive dashboard for mission planners",
        "▸  API for integration with ISRO mission systems",
    ]),
]

for i, (title, accent, items) in enumerate(phases):
    x = Inches(0.3 + i * 4.3)
    add_card_bg(slide, x, Inches(2.0), Inches(4.0), Inches(3.8))
    add_accent_bar(slide, x + Inches(0.1), Inches(2.15),
                   width=Inches(3.8), height=Inches(0.05), color=accent)
    add_text_box(slide, x + Inches(0.2), Inches(2.35), Inches(3.6), Inches(0.5),
                 title, font_size=14, color=accent, bold=True)
    add_multi_text(slide, x + Inches(0.2), Inches(2.9), Inches(3.6), Inches(2.5),
                   items, font_size=14, color=LIGHT)

# Bottom vision
add_card_bg(slide, Inches(2.0), Inches(6.2), Inches(9.3), Inches(1.0),
            color=RGBColor(0x0A, 0x1A, 0x35))
add_text_box(slide, Inches(2.2), Inches(6.3), Inches(8.9), Inches(0.8),
             "VISION: A validated, operational mission planning tool that directly supports\nIndia's next-generation lunar missions — Chandrayaan-4, LUPEX, and beyond",
             font_size=15, color=GOLD, alignment=PP_ALIGN.CENTER, bold=True)

add_bottom_tag(slide)
add_notes(slide, """SPEAKER NOTES — Future Work
Our roadmap has three clear phases.

Phase 1 — Validation in 3 months: Download the missing datasets — calibrated OHRC that actually overlaps our AOI, Diviner thermal for cold-trap confirmation below 110 Kelvin, real PSR illumination products, and complex Stokes SAR products for true CPR/DOP computation. This phase is about replacing proxies with real measurements.

Phase 2 — Enhancement in 6 months: Add multi-pass SAR temporal consistency to see if candidates persist across different observation dates. Integrate LOLA albedo, LAMP ultraviolet water proxy, and M3 spectral data for cross-validation. Enhance rover planning with dynamic illumination and communication window constraints.

Phase 3 — Scaling in 12 months: Extend coverage to 10 or more AOIs across the entire lunar south pole. Build a real-time data ingestion pipeline for ongoing missions. Create an interactive dashboard for ISRO mission planners and an API for integration with mission systems.

Our vision is clear: a validated, operational tool that directly supports India's next-generation lunar missions — Chandrayaan-4, LUPEX, and beyond.""")

# ═══════════════════════════════════════════════════════════════════
#  SLIDE 17 — CONCLUSION
# ═══════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)

# Title
add_text_box(slide, Inches(1.0), Inches(0.8), Inches(11), Inches(0.8),
             "LUNAQUEST", font_size=44, color=ACCENT_1, bold=True,
             alignment=PP_ALIGN.CENTER)
add_text_box(slide, Inches(1.0), Inches(1.6), Inches(11), Inches(0.5),
             "Enabling India's Lunar Ice Discovery Mission",
             font_size=22, color=WHITE, alignment=PP_ALIGN.CENTER)

# Key takeaways
add_card_bg(slide, Inches(1.5), Inches(2.5), Inches(10.3), Inches(3.5))
add_text_box(slide, Inches(1.7), Inches(2.6), Inches(9.8), Inches(0.4),
             "KEY TAKEAWAYS", font_size=20, color=GOLD, bold=True,
             alignment=PP_ALIGN.CENTER)

takeaways = [
    "1️⃣   Complete end-to-end pipeline: Data → Screening → Landing → Rover → Resource Estimation",
    "2️⃣   Built on REAL Chandrayaan-2 data (SAR/DFSAR + TMC-2 DTM, 100% AOI coverage)",
    "3️⃣   Honest about what we know and what we don't — every limitation documented",
    "4️⃣   Roughness-ambiguity addressed head-on — not ignored like most approaches",
    "5️⃣   Fully explainable decisions — 'why this patch, why this site, why this route'",
    "6️⃣   Ready for incremental validation and extension to operational use",
]
add_multi_text(slide, Inches(1.7), Inches(3.1), Inches(9.8), Inches(2.5),
               takeaways, font_size=16, color=LIGHT)

# Thank you
add_text_box(slide, Inches(1.0), Inches(6.3), Inches(11), Inches(0.5),
             "Thank You  —  Team Tantra-Mantra",
             font_size=24, color=GOLD, bold=True, alignment=PP_ALIGN.CENTER)

add_text_box(slide, Inches(1.0), Inches(6.8), Inches(11), Inches(0.4),
             "Bharat Antariksh Hackathon 2026  |  Preliminary candidate screening result; independent validation required.",
             font_size=12, color=LIGHT, alignment=PP_ALIGN.CENTER)

# Bottom accent bar
strip3 = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
         Inches(0), Inches(7.44), SLIDE_W, Inches(0.06))
strip3.fill.solid()
strip3.fill.fore_color.rgb = ACCENT_1
strip3.line.fill.background()

add_notes(slide, """SPEAKER NOTES — Conclusion
To conclude, LunaQuest is not just another hackathon project. It is a complete, scientifically rigorous mission planning framework.

We process real Chandrayaan-2 data — not simulated or synthetic. We achieve 100% AOI coverage validation. We detect, rank, and explain ice candidate patches. We select landing sites, plan rover routes, and estimate resources.

Most importantly, we are honest about what we can and cannot claim. We address the roughness-CPR ambiguity head-on rather than ignoring it. Every decision is explainable. Every limitation is documented.

This framework is ready for incremental validation — as new datasets become available (OHRC, Diviner, PSR), they plug directly into our pipeline and improve confidence.

Thank you judges. We are Team Tantra-Mantra, and we believe LunaQuest can directly support India's next-generation lunar exploration missions.

We're happy to take your questions.""")

# ═══════════════════════════════════════════════════════════════════
#  SAVE
# ═══════════════════════════════════════════════════════════════════
prs.save(OUT_PPT)
print(f"\n[OK] Presentation saved to: {OUT_PPT}")
print(f"     Total slides: {len(prs.slides)}")
